from pptx import Presentation

ppt_path = "c:\\Users\\abhir\\OneDrive\\Documents\\VSCode\\CodeMap\\CodeMap_Phase2_First_Presentation.pptx"
prs = Presentation(ppt_path)

replacements = {
    "File Upload (Planned)": "Project File/Folder Upload (Completed)",
    "File Dependency Graph (Planned)": "Project Blueprint & Dependency Graph (Completed)",
    "Repository Overview (Planned)": "Repository Health Insights (Completed)",
    "Clickable Call Graphs (Planned)": "Interactive Drill-Down & Call Graphs (Completed)",
    "Auto Language Detection (Planned)": "Auto Language Detection (Completed)",
    "Error & Edge Case Handling (Planned)": "Circular Dependency Detection (Completed)",
    "Export Options (Planned)": "Project Analytics Panel (Completed)"
}

for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old_text, new_text in replacements.items():
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

prs.save(ppt_path)
print("Presentation updated successfully.")

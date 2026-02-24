from pptx import Presentation

ppt_path = "c:\\Users\\abhir\\OneDrive\\Documents\\VSCode\\CodeMap\\CodeMap_Phase2_First_Presentation.pptx"
prs = Presentation(ppt_path)

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(f"Shape: {shape.text.encode('utf-8')}")
